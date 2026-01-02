import io
import os
from typing import Optional, Tuple, Any, Dict
import re
import sys
import traceback

from core.services.llm_service import LLMService

class DataAnalysisService:
    def __init__(self, data_manager):
        self.data_manager = data_manager
        self.df: Optional[pd.DataFrame] = None
        self.current_filename: str = ""

    def load_file(self, file_path: str) -> bool:
        """Load a CSV or Excel file into a pandas DataFrame."""
        import pandas as pd
        try:
            from docx import Document
        except ImportError:
            Document = None
        try:
            from pypdf import PdfReader
        except ImportError:
            PdfReader = None
        try:
            from pptx import Presentation as PptxPresentation
        except ImportError:
            PptxPresentation = None

        try:
            filename = os.path.basename(file_path)
            if file_path.endswith('.csv'):
                self.df = pd.read_csv(file_path)
            elif file_path.endswith(('.xls', '.xlsx')):
                self.df = pd.read_excel(file_path)
            elif file_path.endswith('.docx'):
                if Document is None:
                    return False
                doc = Document(file_path)
                data = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        data.append({"Style": para.style.name, "Text": para.text})
                self.df = pd.DataFrame(data)
            elif file_path.endswith('.pdf'):
                if PdfReader is None:
                    return False
                reader = PdfReader(file_path)
                data = []
                for i, page in enumerate(reader.pages):
                    text = page.extract_text()
                    if text:
                        data.append({"Page": i + 1, "Text": text})
                self.df = pd.DataFrame(data)
            elif file_path.endswith('.pptx'):
                if PptxPresentation is None:
                    return False
                prs = PptxPresentation(file_path)
                data = []
                for i, slide in enumerate(prs.slides):
                    text_runs = []
                    title = ""
                    if slide.shapes.title:
                        title = slide.shapes.title.text
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text_runs.append(shape.text)
                    
                    full_text = "\n".join(text_runs)
                    data.append({"Slide": i + 1, "Title": title, "Text": full_text})
                self.df = pd.DataFrame(data)
            else:
                return False
            
            self.current_filename = filename
            return True
        except Exception as e:
            print(f"Error loading file: {e}")
            return False

    def get_basic_stats(self) -> str:
        """Return a string summary of the dataframe."""
        if self.df is None:
            return "Aucune donnée chargée."
        
        buffer = io.StringIO()
        self.df.info(buf=buffer)
        info_str = buffer.getvalue()
        
        stats = f"Nom du fichier: {self.current_filename}\n"
        stats += f"Dimensions: {self.df.shape[0]} lignes, {self.df.shape[1]} colonnes\n\n"
        stats += "Aperçu (5 premières lignes):\n"
        stats += self.df.head().to_string() + "\n\n"
        stats += "Statistiques descriptives:\n"
        stats += self.df.describe().to_string() + "\n\n"
        stats += "Info types:\n"
        stats += info_str
        
        return stats

    def generate_chart(self) -> Optional[Any]:
        """
        Generate a histogram for the first numeric column.
        Returns a matplotlib Figure object or None.
        """
        import matplotlib.pyplot as plt

        if self.df is None:
            return None

        # Find first numeric column
        numeric_cols = self.df.select_dtypes(include=['number']).columns
        if len(numeric_cols) == 0:
            return None
        
        col_name = numeric_cols[0]
        
        # Create figure
        fig, ax = plt.subplots(figsize=(6, 4))
        self.df[col_name].hist(ax=ax, bins=15, color='skyblue', edgecolor='black')
        ax.set_title(f"Distribution de {col_name}")
        ax.set_xlabel(col_name)
        ax.set_ylabel("Fréquence")
        plt.tight_layout()
        
        return fig

    def analyze_with_llm(self, provider_override: Optional[str] = None) -> str:
        """Send data summary to LLM for qualitative analysis."""
        if self.df is None:
            return "Aucune donnée à analyser."

        # Get settings
        settings = self.data_manager.get_settings()
        # Use override, then doc_analyst_provider, then fallback
        provider = provider_override if provider_override else settings.get("doc_analyst_provider", settings.get("chat_provider", "OpenAI GPT-4o mini"))
        api_keys = settings.get("api_keys", {})
        
        # Resolve key
        api_key = api_keys.get(provider)
        if not api_key:
            return f"Clé API manquante pour le provider: {provider}"

        # Prepare Prompt
        stats_summary = self.get_basic_stats()
        # Truncate if too long (simple check)
        if len(stats_summary) > 10000:
            stats_summary = stats_summary[:10000] + "\n...(tronqué)"

        messages = [
            {"role": "system", "content": "Tu es un expert en analyse de données. Analyse le résumé statistique suivant et donne des insights pertinents, des tendances ou des points d'attention. Sois concis et professionnel."},
            {"role": "user", "content": f"Voici les données:\n{stats_summary}"}
        ]

        success, response = LLMService.generate_response(
            provider_name=provider,
            api_key=api_key,
            messages=messages,
            base_url=settings.get("endpoints", {}).get(provider),
            model=settings.get("models", {}).get(provider)
        )

        if success:
            return response
        else:
            if "Endpoint manquant pour IAKA" in response:
                return "Erreur Configuration: L'endpoint IAKA n'est pas configuré.\n\nVeuillez aller dans Administration > Connecteur Chat,\nsélectionner 'IAKA (Interne)' et entrer l'URL de l'endpoint."
            return f"Erreur lors de l'analyse LLM: {response}"

    def generate_code_from_query(self, user_query: str, provider_override: Optional[str] = None) -> Tuple[str, Optional[str]]:
        """
        Step 1: Ask LLM to generate Python code to answer the query.
        Returns (code_str, error_message).
        If success, error_message is None.
        """
        if self.df is None:
            return "", "Veuillez d'abord importer un fichier de données."

        # 1. Get Settings & LLM
        settings = self.data_manager.get_settings()
        provider = provider_override if provider_override else settings.get("doc_analyst_provider", settings.get("chat_provider", "OpenAI GPT-4o mini"))
        api_keys = settings.get("api_keys", {})
        api_key = api_keys.get(provider)
        
        if not api_key:
             return "", f"Clé API manquante pour {provider}"

        # 2. Prepare Prompt
        buffer = io.StringIO()
        self.df.info(buf=buffer)
        info_str = buffer.getvalue()
        
        prompt = f"""
You are a Python Data Analyst. You have a pandas DataFrame named `df` loaded in memory.
Info about `df`:
{info_str}

Columns: {list(self.df.columns)}

User Query: "{user_query}"

Instructions:
1. Write Python code to answer the query.
2. If the user asks for a plot/chart, use `matplotlib.pyplot` (as `plt`).
   - Create a figure explicitly: `fig, ax = plt.subplots()`
   - Plot on `ax`.
   - IMPORTANT: Do NOT call `plt.show()`.
   - The figure object must be available as `fig`.
3. If the user asks for a calculation, print the result using `print()`.
4. Wrap your code in a markdown block: ```python ... ```.
5. Do NOT provide explanations outside the code block, just the code.
6. Keep it simple and robust. Handle potential NaNs if necessary.
"""
        
        messages = [{"role": "user", "content": prompt}]
        
        success, response = LLMService.generate_response(
            provider_name=provider,
            api_key=api_key,
            messages=messages,
            base_url=settings.get("endpoints", {}).get(provider),
            model=settings.get("models", {}).get(provider)
        )
        
        if not success:
            return "", f"Erreur IA: {response}"

        # 3. Extract Code
        code_match = re.search(r"```python(.*?)```", response, re.DOTALL)
        if code_match:
            code = code_match.group(1).strip()
        else:
            if "def " in response or "print(" in response or "fig" in response:
                code = response
            else:
                return "", f"L'IA n'a pas généré de code valide.\nRéponse: {response}"
                
        return code, None

    def execute_generated_code(self, code: str) -> Tuple[str, Optional[Any]]:
        """
        Step 2: Execute the provided python code.
        WARNING: This executes arbitrary code. Ensure user validation before calling.
        """
        if self.df is None:
            return "Erreur: DataFrame non chargé.", None
            
        import pandas as pd
        import matplotlib.pyplot as plt

        # We need to capture stdout
        old_stdout = sys.stdout
        redirected_output = io.StringIO()
        sys.stdout = redirected_output
        
        # Context for execution
        local_scope = {"df": self.df, "pd": pd, "plt": plt}
        fig = None
        
        try:
            exec(code, {}, local_scope)
            
            # Check if 'fig' variable exists in local_scope
            if 'fig' in local_scope and isinstance(local_scope['fig'], plt.Figure):
                fig = local_scope['fig']
                
        except Exception as e:
            sys.stdout = old_stdout
            return f"Erreur d'exécution du code généré:\n{e}\n\nCode:\n{code}", None
            
        sys.stdout = old_stdout
        output_text = redirected_output.getvalue()
        
        if not output_text and not fig:
            output_text = "Code exécuté avec succès, mais aucun résultat affiché."
            
        return output_text, fig

    def export_to_pptx(self, output_path: str, llm_analysis: str = "") -> bool:
        """Export analysis to a PowerPoint presentation."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            import matplotlib.pyplot as plt

            prs = Presentation()

            # Slide 1: Title
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = f"Analyse de {self.current_filename}"
            subtitle.text = "Généré par Mon Assistant Perso"

            # Slide 2: Stats Summary
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            title_shape.text = "Résumé Statistique"
            
            tf = body_shape.text_frame
            tf.text = "Aperçu des données"
            
            # Simple stats integration
            if self.df is not None:
                p = tf.add_paragraph()
                p.text = f"Lignes: {self.df.shape[0]}"
                p = tf.add_paragraph()
                p.text = f"Colonnes: {self.df.shape[1]}"
                
                # Add columns list
                p = tf.add_paragraph()
                p.text = "Colonnes: " + ", ".join(self.df.columns[:5]) + ("..." if len(self.df.columns) > 5 else "")

            # Slide 3: LLM Analysis
            if llm_analysis:
                slide = prs.slides.add_slide(bullet_slide_layout)
                title_shape = slide.shapes.title
                title_shape.text = "Analyse IA"
                body_shape = slide.shapes.placeholders[1]
                body_shape.text_frame.text = llm_analysis[:1000] # Truncate to fit roughly

            # Slide 4: Chart
            fig = self.generate_chart()
            if fig:
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                
                image_stream = io.BytesIO()
                fig.savefig(image_stream, format='png')
                image_stream.seek(0)
                
                left = Inches(1)
                top = Inches(1.5)
                height = Inches(4.5)
                slide.shapes.add_picture(image_stream, left, top, height=height)
                
                title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
                title_box.text_frame.text = "Visualisation (Premier champ numérique)"
                
                # Close figure to free memory
                plt.close(fig)

            prs.save(output_path)
            return True
        except Exception as e:
            print(f"Error exporting to PPTX: {e}")
            import traceback
            traceback.print_exc()
            return False
